from __future__ import annotations

import logging
from pathlib import Path

from sqlalchemy import update
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.config import get_settings
from app.models.knowledge import KnowledgeChunk
from app.models.reference_file import ReferenceFile
from app.models.transcript import TranscriptSegment
from app.services.embedding import EmbeddingService
from app.services.semantic_pipeline import RawChunk, SemanticPipelineService
from app.services.storage import StorageService


settings = get_settings()
logger = logging.getLogger("app.reference")


class ReferenceProcessingService:
    def __init__(self) -> None:
        self.embedding_service = EmbeddingService()
        self.semantic_pipeline = SemanticPipelineService()
        self.storage_service = StorageService()
        self.match_threshold = settings.reference_match_threshold

    async def process_reference_files(
        self,
        session: AsyncSession,
        lecture_id,
        reference_files: list[ReferenceFile],
        transcript_segments: list[TranscriptSegment],
    ) -> dict:
        if not reference_files:
            return {
                "reference_file_count": 0,
                "reference_files_matched": 0,
                "reference_files_skipped": 0,
                "reference_chunks": 0,
                "reference_status": "none",
            }

        lecture_text = " ".join((segment.edited_text or segment.text) for segment in transcript_segments).strip()
        if not lecture_text:
            return {
                "reference_file_count": len(reference_files),
                "reference_files_matched": 0,
                "reference_files_skipped": len(reference_files),
                "reference_chunks": 0,
                "reference_status": "missing_lecture_text",
            }

        lecture_embedding = self.embedding_service.encode([lecture_text])[0]
        chunk_count = 0
        matched = 0
        skipped = 0

        for reference_file in reference_files:
            try:
                resolved_path = self.storage_service.ensure_local_path(reference_file.storage_path, reference_file.details)
                reference_text = self.extract_text(resolved_path, reference_file.file_type)
            except Exception as exc:
                logger.warning(
                    "reference_extract_failed lecture=%s file=%s error=%s",
                    lecture_id,
                    reference_file.original_filename,
                    exc,
                )
                reference_file.details = {
                    **(reference_file.details or {}),
                    "processing_status": "extract_failed",
                    "processing_error": str(exc),
                }
                skipped += 1
                continue

            if not reference_text.strip():
                reference_file.details = {
                    **(reference_file.details or {}),
                    "processing_status": "empty_text",
                }
                skipped += 1
                continue

            reference_embedding = self.embedding_service.encode([reference_text])[0]
            similarity = self.cosine_similarity(lecture_embedding, reference_embedding)
            if similarity < self.match_threshold:
                logger.info(
                    "reference_skipped_mismatch lecture=%s file=%s similarity=%s threshold=%s",
                    lecture_id,
                    reference_file.original_filename,
                    similarity,
                    self.match_threshold,
                )
                reference_file.details = {
                    **(reference_file.details or {}),
                    "processing_status": "skipped_mismatch",
                    "match_score": similarity,
                    "match_threshold": self.match_threshold,
                }
                skipped += 1
                continue

            chunks = [RawChunk(text=reference_text, start_time=0.0, end_time=0.0)]
            _cleaned_text, sentence_units, topic_units = self.semantic_pipeline.build_from_chunks(chunks)
            content_blocks: list[tuple[str, str, dict]] = []
            content_blocks.append(
                (
                    f"Reference {reference_file.original_filename}",
                    reference_text,
                    {
                        "kind": "reference_full",
                        "reference_file_id": str(reference_file.id),
                        "reference_filename": reference_file.original_filename,
                        "match_score": similarity,
                        "trusted_source": True,
                        "approved_for_kb": True,
                        "student_visible": True,
                    },
                )
            )
            for topic in topic_units:
                content_blocks.append(
                    (
                        topic.title,
                        topic.content or topic.description,
                        {
                            "kind": "reference_topic",
                            "reference_file_id": str(reference_file.id),
                            "reference_filename": reference_file.original_filename,
                            "topic_sequence": topic.sequence,
                            "summary": topic.description,
                            "match_score": similarity,
                            "trusted_source": True,
                            "approved_for_kb": True,
                            "student_visible": True,
                        },
                    )
                )

            embeddings = self.embedding_service.encode([content for _, content, _ in content_blocks]) if content_blocks else []
            for (topic, content, metadata), embedding in zip(content_blocks, embeddings, strict=True):
                session.add(
                    KnowledgeChunk(
                        lecture_id=lecture_id,
                        topic=topic,
                        content=content,
                        details=metadata,
                        embedding=embedding,
                    )
                )
                chunk_count += 1

            reference_file.details = {
                **(reference_file.details or {}),
                "processing_status": "matched",
                "match_score": similarity,
                "match_threshold": self.match_threshold,
                "sentence_count": len(sentence_units),
                "topic_count": len(topic_units),
            }
            matched += 1

        await session.flush()
        return {
            "reference_file_count": len(reference_files),
            "reference_files_matched": matched,
            "reference_files_skipped": skipped,
            "reference_chunks": chunk_count,
            "reference_status": "matched" if matched else "skipped",
        }

    def extract_text(self, storage_path: str, file_type: str) -> str:
        normalized = file_type.lower()
        path = Path(storage_path)
        if normalized == "pdf":
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            return "\n".join((page.extract_text() or "").strip() for page in reader.pages).strip()
        if normalized in {"ppt", "pptx"}:
            from pptx import Presentation

            presentation = Presentation(str(path))
            slides: list[str] = []
            for slide in presentation.slides:
                texts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                if texts:
                    slides.append("\n".join(texts))
            return "\n\n".join(slides).strip()
        raise ValueError(f"Unsupported reference file type: {file_type}")

    def cosine_similarity(self, left: list[float], right: list[float]) -> float:
        if not left or not right or len(left) != len(right):
            return 0.0
        return round(sum(a * b for a, b in zip(left, right, strict=True)), 4)
